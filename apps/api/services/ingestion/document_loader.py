"""Unified content extraction — Crawl4AI as primary engine + GPT-Researcher loader_dict for Office formats.

Architecture:
- Crawl4AI: web pages (HTTP/HTTPS), PDF (URL + local via file://), local HTML
- loader_dict: DOCX, PPTX, XLSX, CSV, TXT, MD (Office/text formats Crawl4AI doesn't support)
- Fallback layers: httpx + clean_soup, trafilatura, Playwright cascade

Features:
- BM25 content filtering for focused extraction from noisy web pages
- Batch URL crawling via arun_many() for bulk ingestion
- Media/links metadata extraction and storage

References:
- Crawl4AI: AsyncWebCrawler, NaivePDFProcessorStrategy (processors/pdf/processor.py)
- Crawl4AI: BM25ContentFilter (content_filter_strategy.py L381-530)
- GPT-Researcher: DocumentLoader loader_dict (document/document.py L66-78)
- GPT-Researcher: clean_soup(), get_text_from_soup() (scraper/utils.py L94-132)
"""

import logging
import re
import asyncio
from pathlib import Path
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)

# Office formats handled by loader_dict (GPT-Researcher pattern)
OFFICE_EXTENSIONS = {"doc", "docx", "pptx", "csv", "xls", "xlsx"}

# Text formats read directly
TEXT_EXTENSIONS = {"txt", "md", "rst"}
_marker_models: dict | None = None


@dataclass
class ExtractionResult:
    """Rich extraction result with metadata from Crawl4AI."""

    title: str = ""
    content: str = ""
    images: list[dict] = field(default_factory=list)    # [{url, score, alt, ...}]
    links_internal: list[dict] = field(default_factory=list)
    links_external: list[dict] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)         # OG/Twitter/standard meta tags


async def extract_content(
    file_path: str | None = None,
    url: str | None = None,
    query: str | None = None,
    session_name: str | None = None,
) -> tuple[str, str]:
    """Unified content extraction entry point.

    Returns (title, markdown_content).
    Routes to the appropriate extractor based on input type and format.

    Args:
        query: Optional search query for BM25 content filtering (Crawl4AI).
               When provided, noisy web pages are filtered to keep only
               query-relevant sections.
        session_name: Optional Playwright session name for authenticated Canvas API access.
    """
    if url:
        return await _extract_from_url(url, query=query, session_name=session_name)

    if not file_path:
        return "", ""

    ext = Path(file_path).suffix.lstrip(".").lower()

    # PDF and HTML → Crawl4AI via file:// protocol
    if ext == "pdf":
        return await _extract_local_with_crawl4ai(file_path)
    if ext in ("html", "htm"):
        return await _extract_local_with_crawl4ai(file_path)

    # Office formats → GPT-Researcher loader_dict
    if ext in OFFICE_EXTENSIONS:
        return await asyncio.to_thread(_extract_with_loader_dict, file_path, ext)

    # Text formats → direct read
    if ext in TEXT_EXTENSIONS:
        return await asyncio.to_thread(_extract_plain_text, file_path)

    # Unknown → try text read
    return await asyncio.to_thread(_extract_plain_text, file_path)


def _build_crawl4ai_config(query: str | None = None):
    """Build Crawl4AI CrawlerRunConfig with optional BM25 filtering."""
    from crawl4ai import CrawlerRunConfig
    from crawl4ai.markdown_generation_strategy import DefaultMarkdownGenerator

    markdown_generator = DefaultMarkdownGenerator()
    if query:
        try:
            from crawl4ai.content_filter_strategy import BM25ContentFilter

            markdown_generator = DefaultMarkdownGenerator(
                content_filter=BM25ContentFilter(user_query=query, bm25_threshold=1.0)
            )
        except ImportError:
            pass
    return CrawlerRunConfig(
        excluded_tags=["nav", "footer", "sidebar"],
        word_count_threshold=100,
        markdown_generator=markdown_generator,
    )


async def extract_content_rich(
    url: str,
    query: str | None = None,
) -> ExtractionResult:
    """Rich URL extraction returning content with media, links, and metadata.

    Uses Crawl4AI with optional BM25 filtering. Falls back to basic extraction.
    """
    try:
        from crawl4ai import AsyncWebCrawler

        async with AsyncWebCrawler() as crawler:
            config = _build_crawl4ai_config(query)
            result = await crawler.arun(url=url, config=config)
            if result.success:
                return _crawl_result_to_extraction(result, url, query)
    except ImportError:
        logger.debug("crawl4ai not installed, falling back to basic extraction")
    except Exception as e:
        logger.debug(f"Crawl4AI rich extraction failed for {url}: {e}")

    # Fallback to basic extraction
    title, content = await _extract_from_url(url, query=query)
    return ExtractionResult(title=title, content=content)


async def extract_content_batch(
    urls: list[str],
    query: str | None = None,
) -> list[ExtractionResult]:
    """Batch URL extraction using Crawl4AI arun_many() for parallel crawling.

    Significantly faster than sequential extraction for multiple URLs.
    Falls back to sequential extraction if Crawl4AI is unavailable.
    """
    if not urls:
        return []

    try:
        from crawl4ai import AsyncWebCrawler

        config = _build_crawl4ai_config(query)

        async with AsyncWebCrawler() as crawler:
            results = await crawler.arun_many(urls=urls, config=config)
            extractions = []
            for result in results:
                if result.success:
                    extractions.append(
                        _crawl_result_to_extraction(result, result.url, query)
                    )
                else:
                    extractions.append(ExtractionResult(title=result.url))
            return extractions

    except ImportError:
        logger.debug("crawl4ai not installed, falling back to sequential extraction")
    except Exception as e:
        logger.debug(f"Crawl4AI batch extraction failed: {e}")

    # Fallback: sequential extraction
    results = []
    for url in urls:
        title, content = await _extract_from_url(url, query=query)
        results.append(ExtractionResult(title=title, content=content))
    return results


def _crawl_result_to_extraction(result, url: str, query: str | None = None) -> ExtractionResult:
    """Convert a Crawl4AI CrawlResult to ExtractionResult with media/links/metadata."""
    md = getattr(result, "markdown", None)
    if md is None:
        return ExtractionResult(title=url)

    # Content: prefer BM25-filtered when query was provided
    content = ""
    if query and hasattr(md, "fit_markdown") and md.fit_markdown:
        content = md.fit_markdown
    if not content:
        content = md.raw_markdown if hasattr(md, "raw_markdown") else str(md)

    # Title from metadata
    title = url
    if result.metadata and isinstance(result.metadata, dict):
        title = result.metadata.get("title", url) or url

    # Images — Crawl4AI provides scored image data
    images = []
    if hasattr(result, "media") and isinstance(result.media, dict):
        for img in result.media.get("images", []):
            images.append({
                "src": img.get("src", ""),
                "alt": img.get("alt", ""),
                "score": img.get("score", 0),
                "description": img.get("desc", ""),
            })

    # Links — internal and external
    links_internal = []
    links_external = []
    if hasattr(result, "links") and isinstance(result.links, dict):
        for link in result.links.get("internal", []):
            links_internal.append({
                "href": link.get("href", ""),
                "text": link.get("text", ""),
            })
        for link in result.links.get("external", []):
            links_external.append({
                "href": link.get("href", ""),
                "text": link.get("text", ""),
            })

    # Metadata
    metadata = {}
    if result.metadata and isinstance(result.metadata, dict):
        metadata = result.metadata

    return ExtractionResult(
        title=title,
        content=content,
        images=images,
        links_internal=links_internal,
        links_external=links_external,
        metadata=metadata,
    )


# ── Canvas API extraction (moved to canvas_loader.py) ──
# Import Canvas functions for use by _extract_from_url and other local callers
from services.ingestion.canvas_loader import (  # noqa: E402
    _try_canvas_api,
    CanvasAuthExpiredError,
)


# ── Crawl4AI extraction (web + PDF + HTML) ──


async def _extract_from_url(
    url: str,
    query: str | None = None,
    session_name: str | None = None,
) -> tuple[str, str]:
    """URL extraction — multi-layer cascade with Canvas API priority.

    Fallback cascade:
    0. Canvas REST API (structured data, no HTML scraping needed)
    1. Crawl4AI (best quality — Markdown + media + metadata + BM25 filtering)
    2. httpx + clean_soup (GPT-Researcher pattern)
    3. trafilatura
    4. Playwright (existing cascade from automation.py)
    """
    # Layer 0: Canvas REST API (structured data — bypasses HTML scraping)
    result = await _try_canvas_api(url, session_name=session_name)
    if result:
        return result

    # Layer 1: Crawl4AI (with optional BM25 content filtering)
    result = await _try_crawl4ai_url(url, query=query)
    if result:
        return result

    # Layer 2: httpx + clean_soup (GPT-Researcher BeautifulSoupScraper pattern)
    result = await _try_httpx_clean_soup(url)
    if result:
        return result

    # Layer 3: trafilatura
    result = await _try_trafilatura_url(url)
    if result:
        return result

    # Layer 4: Playwright (existing browser cascade)
    result = await _try_browser_cascade(url)
    if result:
        return result

    return "", ""


async def _try_crawl4ai_url(url: str, query: str | None = None) -> tuple[str, str] | None:
    """Layer 1: Crawl4AI — handles web pages and PDF URLs uniformly.

    When query is provided, enables BM25 content filtering to extract only
    query-relevant sections from noisy web pages (Crawl4AI BM25ContentFilter).
    """
    try:
        from crawl4ai import AsyncWebCrawler

        async with AsyncWebCrawler() as crawler:
            config = _build_crawl4ai_config(query)
            result = await crawler.arun(url=url, config=config)
            if result.success:
                md = result.markdown
                # Prefer BM25-filtered markdown when query was provided
                raw = ""
                if query and hasattr(md, "fit_markdown") and md.fit_markdown:
                    raw = md.fit_markdown
                if not raw:
                    raw = md.raw_markdown if hasattr(md, "raw_markdown") else str(md)
                if len(raw) >= 100:
                    title = url
                    if result.metadata and isinstance(result.metadata, dict):
                        title = result.metadata.get("title", url) or url
                    return title, raw
    except ImportError:
        logger.debug("crawl4ai not installed, skipping Crawl4AI layer")
    except Exception as e:
        logger.debug(f"Crawl4AI failed for {url}: {e}")
    return None


async def _try_httpx_clean_soup(url: str) -> tuple[str, str] | None:
    """Layer 2: httpx fetch + GPT-Researcher clean_soup() HTML cleaning."""
    try:
        import httpx
        from bs4 import BeautifulSoup

        async with httpx.AsyncClient(follow_redirects=True, timeout=15) as client:
            resp = await client.get(url)
            if resp.status_code != 200:
                return None
            html = resp.text

        soup = BeautifulSoup(html, "lxml")
        soup = clean_soup(soup)
        title = extract_title(soup)
        content = get_text_from_soup(soup)

        if len(content) >= 100:
            return title or url, content
    except ImportError:
        logger.debug("bs4/lxml not installed, skipping httpx+clean_soup layer")
    except Exception as e:
        logger.debug(f"httpx+clean_soup failed for {url}: {e}")
    return None


async def _try_trafilatura_url(url: str) -> tuple[str, str] | None:
    """Layer 3: trafilatura fallback (runs in thread to avoid blocking loop)."""
    return await asyncio.to_thread(_try_trafilatura_url_sync, url)


def _try_trafilatura_url_sync(url: str) -> tuple[str, str] | None:
    """Synchronous trafilatura fallback implementation."""
    try:
        import trafilatura

        downloaded = trafilatura.fetch_url(url)
        if not downloaded:
            return None
        content = trafilatura.extract(
            downloaded,
            include_links=True,
            include_formatting=True,
            include_tables=True,
            output_format="txt",
        )
        if content and len(content) >= 100:
            metadata = trafilatura.extract_metadata(downloaded)
            title = metadata.title if metadata and metadata.title else url
            return title, content
    except ImportError:
        logger.debug("trafilatura not installed, skipping trafilatura layer")
    except Exception as e:
        logger.debug(f"trafilatura failed for {url}: {e}")
    return None


async def _try_browser_cascade(url: str) -> tuple[str, str] | None:
    """Layer 4: Playwright browser cascade (existing automation.py)."""
    try:
        from services.browser.automation import cascade_fetch
        from services.parser.url import extract_text_from_html

        html = await cascade_fetch(url)
        if html:
            text = extract_text_from_html(html)
            if text and len(text) >= 50:
                return url, text
    except Exception as e:
        logger.debug(f"Browser cascade failed for {url}: {e}")
    return None


# ── Crawl4AI local file extraction ──


async def _extract_local_with_crawl4ai(file_path: str) -> tuple[str, str]:
    """Extract content from local PDF/HTML via Crawl4AI file:// protocol.

    Falls back to legacy extractors if Crawl4AI is unavailable.
    """
    abs_path = str(Path(file_path).resolve())
    file_url = f"file://{abs_path}"

    try:
        from crawl4ai import AsyncWebCrawler, CrawlerRunConfig

        async with AsyncWebCrawler() as crawler:
            result = await crawler.arun(url=file_url, config=CrawlerRunConfig())
            if result.success:
                md = result.markdown
                raw = md.raw_markdown if hasattr(md, "raw_markdown") else str(md)
                if raw:
                    title = ""
                    if result.metadata and isinstance(result.metadata, dict):
                        title = result.metadata.get("title", "") or ""
                    return title or Path(file_path).stem, raw
    except ImportError:
        logger.debug("crawl4ai not installed, falling back to legacy extractors")
    except Exception as e:
        logger.debug(f"Crawl4AI local extraction failed for {file_path}: {e}")

    # Fallback for PDF
    ext = Path(file_path).suffix.lstrip(".").lower()
    if ext == "pdf":
        return await asyncio.to_thread(_extract_pdf_fallback, file_path)

    # Fallback for HTML
    if ext in ("html", "htm"):
        return await asyncio.to_thread(_extract_html_fallback, file_path)

    return Path(file_path).stem, ""


def _extract_pdf_fallback(file_path: str) -> tuple[str, str]:
    """PDF fallback: Marker → pypdf."""
    # Try Marker
    try:
        from marker.converters.pdf import PdfConverter

        converter = PdfConverter(artifact_dict=_get_marker_models())
        rendered = converter(file_path)
        return Path(file_path).stem, rendered.markdown
    except ImportError:
        pass
    except Exception as e:
        logger.debug(f"Marker PDF extraction failed: {e}")

    # Try pypdf
    try:
        import pypdf

        reader = pypdf.PdfReader(file_path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return Path(file_path).stem, text
    except Exception as e:
        logger.debug(f"pypdf extraction failed: {e}")

    return Path(file_path).stem, ""


def _extract_html_fallback(file_path: str) -> tuple[str, str]:
    """HTML fallback: trafilatura → raw read."""
    try:
        import trafilatura

        with open(file_path) as f:
            html = f.read()
        content = trafilatura.extract(html, include_tables=True) or ""
        if content:
            return Path(file_path).stem, content
    except Exception as e:
        logger.debug("trafilatura HTML extraction failed for %s: %s", file_path, e)

    # Raw text fallback
    try:
        return Path(file_path).stem, Path(file_path).read_text(errors="ignore")
    except Exception as e:
        logger.warning("Failed to read file %s: %s", file_path, e)
        return Path(file_path).stem, ""


# ── GPT-Researcher loader_dict (Office formats) ──


def _extract_with_loader_dict(file_path: str, ext: str) -> tuple[str, str]:
    """DOCX/PPTX/XLSX/CSV extraction using GPT-Researcher's loader_dict pattern.

    Reference: gpt_researcher/document/document.py L66-78
    """
    try:
        from langchain_community.document_loaders import (
            UnstructuredCSVLoader,
            UnstructuredExcelLoader,
            UnstructuredPowerPointLoader,
            UnstructuredWordDocumentLoader,
        )

        loader_map = {
            "doc": lambda fp: UnstructuredWordDocumentLoader(fp),
            "docx": lambda fp: UnstructuredWordDocumentLoader(fp),
            "pptx": lambda fp: UnstructuredPowerPointLoader(fp),
            "csv": lambda fp: UnstructuredCSVLoader(fp, mode="elements"),
            "xls": lambda fp: UnstructuredExcelLoader(fp, mode="elements"),
            "xlsx": lambda fp: UnstructuredExcelLoader(fp, mode="elements"),
        }

        factory = loader_map.get(ext)
        if not factory:
            return Path(file_path).stem, ""

        loader = factory(file_path)
        docs = loader.load()
        content = "\n\n".join(doc.page_content for doc in docs if doc.page_content)
        return Path(file_path).stem, content

    except ImportError:
        logger.warning(
            "langchain-community/unstructured not installed. "
            "Falling back to basic extraction for %s files. "
            "Install with: pip install langchain-community unstructured",
            ext,
        )
        return _extract_office_fallback(file_path, ext)
    except Exception as e:
        logger.warning(f"loader_dict extraction failed for {file_path}: {e}")
        return _extract_office_fallback(file_path, ext)


def _extract_office_fallback(file_path: str, ext: str) -> tuple[str, str]:
    """Basic fallback for Office formats when Unstructured is unavailable."""
    if ext in ("doc", "docx"):
        try:
            from docx import Document

            doc = Document(file_path)
            content = "\n\n".join(p.text for p in doc.paragraphs if p.text)
            return Path(file_path).stem, content
        except Exception as e:
            logger.warning("DOCX fallback extraction failed for %s: %s", file_path, e)

    if ext in ("pptx",):
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            texts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                texts.append(f"## Slide {slide_num}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return Path(file_path).stem, "\n\n".join(texts)
        except Exception as e:
            logger.warning("PPTX fallback extraction failed for %s: %s", file_path, e)

    return Path(file_path).stem, ""


# ── Plain text extraction ──


def _extract_plain_text(file_path: str) -> tuple[str, str]:
    """Direct text read for .txt, .md, .rst, and unknown files."""
    try:
        content = Path(file_path).read_text(errors="ignore")
        return Path(file_path).stem, content
    except Exception as e:
        logger.warning("Failed to read plain text file %s: %s", file_path, e)
        return Path(file_path).stem, ""


def _get_marker_models() -> dict:
    """Lazy-load Marker models once for local PDF fallback path."""
    global _marker_models
    if _marker_models is None:
        from marker.models import create_model_dict

        _marker_models = create_model_dict()
    return _marker_models


# ── HTML cleaning utilities (GPT-Researcher scraper/utils.py L94-132) ──


def clean_soup(soup):
    """Clean HTML by removing unwanted tags.

    Ported from GPT-Researcher scraper/utils.py clean_soup().
    """
    import bs4

    for tag in soup.find_all(
        ["script", "style", "footer", "header", "nav", "menu", "sidebar", "svg"]
    ):
        tag.decompose()

    disallowed_class_set = {"nav", "menu", "sidebar", "footer"}

    def has_disallowed_class(elem):
        if not isinstance(elem, bs4.Tag):
            return False
        return any(cls in disallowed_class_set for cls in elem.get("class", []))

    for tag in soup.find_all(has_disallowed_class):
        tag.decompose()

    return soup


def clean_soup_canvas_aware(soup):
    """Canvas-aware HTML cleaner that preserves course content containers.

    Canvas LMS puts content inside elements that generic cleaners strip out
    (e.g. nav-like sidebars, module containers). This cleaner only removes
    truly non-content elements while preserving Canvas-specific structure.
    """
    # Only remove script, style, and SVG — preserve everything else
    for tag in soup.find_all(["script", "style", "svg", "noscript"]):
        tag.decompose()

    # Remove Canvas chrome (global nav, breadcrumbs, footer) but keep content
    canvas_chrome_ids = [
        "header", "menu", "left-side", "breadcrumbs",
        "flash_message_holder", "footer",
    ]
    for chrome_id in canvas_chrome_ids:
        elem = soup.find(id=chrome_id)
        if elem:
            elem.decompose()

    # Try to extract just the main content area if it exists
    content_area = (
        soup.find(id="content")
        or soup.find(id="wiki_page_show")
        or soup.find(class_="ic-app-main-content")
        or soup.find(role="main")
    )
    if content_area:
        return content_area

    return soup


def extract_title(soup) -> str:
    """Extract title from BeautifulSoup object.

    Ported from GPT-Researcher scraper/utils.py extract_title().
    """
    return soup.title.string if soup.title else ""


def get_text_from_soup(soup) -> str:
    """Get clean text from BeautifulSoup object.

    Ported from GPT-Researcher scraper/utils.py get_text_from_soup().
    """
    text = soup.get_text(strip=True, separator="\n")
    # Collapse runs of 3+ newlines to double newline (paragraph break)
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Collapse runs of spaces (but preserve newlines for structure)
    text = re.sub(r"[^\S\n]{2,}", " ", text)
    return text.strip()


# Backward-compatible re-exports from extracted modules
from services.ingestion.canvas_loader import (  # noqa: E402, F401
    _load_session_cookies,
    CanvasAuthExpiredError,
    CanvasExtraction,
    _extract_file_urls_from_html,
    _canvas_api_paginate,
    _try_canvas_api,
    _try_canvas_api_deep,
    download_canvas_file,
    _parse_canvas_quiz_question,
)
