"""Local file format extraction for document_loader.

Contains:
- Local PDF/HTML extraction via Crawl4AI (with fallbacks)
- Office format extraction via GPT-Researcher loader_dict
- Dedicated fallback extractors for DOCX, PPTX, XLSX, CSV
- Plain text extraction
- Magic-bytes routing for misnamed files (e.g. Canvas .pdf that is actually .xlsx)
- Marker model loading

Extracted from document_loader.py.
"""

import csv
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Office formats handled by loader_dict (GPT-Researcher pattern)
OFFICE_EXTENSIONS = {"doc", "docx", "pptx", "csv", "xls", "xlsx"}

# Text formats read directly
TEXT_EXTENSIONS = {"txt", "md", "rst"}

_marker_models: dict | None = None


async def _extract_local_with_crawl4ai(file_path: str) -> tuple[str, str]:
    """Extract content from local PDF/HTML via Crawl4AI file:// protocol.

    Falls back to legacy extractors if Crawl4AI is unavailable.
    """
    import asyncio

    abs_path = str(Path(file_path).resolve())
    file_url = f"file://{abs_path}"

    try:
        from crawl4ai import AsyncWebCrawler, CrawlerRunConfig

        async with AsyncWebCrawler() as crawler:
            result = await crawler.arun(url=file_url, config=CrawlerRunConfig())
            if result.success:
                md = result.markdown
                raw = md.raw_markdown if hasattr(md, "raw_markdown") else str(md)
                if raw:
                    title = ""
                    if result.metadata and isinstance(result.metadata, dict):
                        title = result.metadata.get("title", "") or ""
                    return title or Path(file_path).stem, raw
    except ImportError:
        logger.debug("crawl4ai not installed, falling back to legacy extractors")
    except (IOError, OSError) as e:
        logger.debug(f"Crawl4AI local file I/O error for {file_path}: {e}")
    except (ValueError, RuntimeError) as e:
        logger.exception(f"Unexpected error in Crawl4AI local extraction for {file_path}")

    # Fallback based on file extension
    ext = Path(file_path).suffix.lstrip(".").lower()
    if ext == "pdf":
        return await asyncio.to_thread(_extract_pdf_fallback, file_path)

    if ext in ("html", "htm"):
        return await asyncio.to_thread(_extract_html_fallback, file_path)

    if ext in ("docx", "doc"):
        return await asyncio.to_thread(_extract_docx_fallback, file_path)

    if ext == "pptx":
        return await asyncio.to_thread(_extract_pptx_fallback, file_path)

    if ext == "xlsx":
        return await asyncio.to_thread(_extract_xlsx_fallback, file_path)

    if ext == "csv":
        return await asyncio.to_thread(_extract_csv_fallback, file_path)

    if ext in TEXT_EXTENSIONS:
        return await asyncio.to_thread(_extract_plain_text, file_path)

    return Path(file_path).stem, ""


def _extract_pdf_fallback(file_path: str) -> tuple[str, str]:
    """PDF fallback: Marker -> pypdf."""
    # Guard: skip files that aren't actually PDFs (e.g. .xlsx/.zip with .pdf extension)
    # PK magic bytes (0x504B) indicate a ZIP-based Office format (docx/pptx/xlsx)
    try:
        with open(file_path, "rb") as f:
            magic = f.read(4)
        if magic[:2] == b"PK":
            logger.info(
                "File %s has .pdf extension but PK magic bytes — attempting Office extraction",
                Path(file_path).name,
            )
            return _extract_pk_office_file(file_path)
        if magic != b"%PDF":
            logger.warning("Skipping non-PDF file disguised as PDF: %s (magic: %r)", Path(file_path).name, magic)
            return Path(file_path).stem, ""
    except OSError:
        pass

    # Try Marker
    try:
        from marker.converters.pdf import PdfConverter

        converter = PdfConverter(artifact_dict=_get_marker_models())
        rendered = converter(file_path)
        return Path(file_path).stem, rendered.markdown
    except ImportError:
        pass
    except (IOError, OSError) as e:
        logger.warning("Marker PDF file I/O error: %s", e)
    except (ValueError, RuntimeError, Exception) as e:
        logger.warning("Marker PDF parsing error for %s: %s", Path(file_path).name, e)

    # Try pypdf
    try:
        import pypdf
        from pypdf.errors import PyPdfError

        reader = pypdf.PdfReader(file_path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return Path(file_path).stem, text
    except ImportError:
        logger.debug("pypdf not installed, skipping PDF fallback")
    except (IOError, OSError) as e:
        logger.warning("pypdf file I/O error: %s", e)
    except (PyPdfError, ValueError, KeyError) as e:
        logger.warning("pypdf parsing error for %s: %s", Path(file_path).name, e)

    return Path(file_path).stem, ""


def _extract_html_fallback(file_path: str) -> tuple[str, str]:
    """HTML fallback: trafilatura -> raw read."""
    try:
        import trafilatura

        with open(file_path) as f:
            html = f.read()
        content = trafilatura.extract(html, include_tables=True) or ""
        if content:
            return Path(file_path).stem, content
    except ImportError:
        logger.debug("trafilatura not installed for HTML extraction of %s", file_path)
    except (IOError, OSError) as e:
        logger.debug("trafilatura HTML file I/O error for %s: %s", file_path, e)
    except (ValueError, KeyError) as e:
        logger.debug("trafilatura HTML parsing error for %s: %s", file_path, e)

    # Raw text fallback
    try:
        return Path(file_path).stem, Path(file_path).read_text(errors="ignore")
    except (IOError, OSError) as e:
        logger.exception("Failed to read file %s", file_path)
        return Path(file_path).stem, ""


def _extract_with_loader_dict(file_path: str, ext: str) -> tuple[str, str]:
    """DOCX/PPTX/XLSX/CSV extraction using GPT-Researcher's loader_dict pattern.

    Reference: gpt_researcher/document/document.py L66-78
    """
    try:
        from langchain_community.document_loaders import (
            UnstructuredCSVLoader,
            UnstructuredExcelLoader,
            UnstructuredPowerPointLoader,
            UnstructuredWordDocumentLoader,
        )

        loader_map = {
            "doc": lambda fp: UnstructuredWordDocumentLoader(fp),
            "docx": lambda fp: UnstructuredWordDocumentLoader(fp),
            "pptx": lambda fp: UnstructuredPowerPointLoader(fp),
            "csv": lambda fp: UnstructuredCSVLoader(fp, mode="elements"),
            "xls": lambda fp: UnstructuredExcelLoader(fp, mode="elements"),
            "xlsx": lambda fp: UnstructuredExcelLoader(fp, mode="elements"),
        }

        factory = loader_map.get(ext)
        if not factory:
            return Path(file_path).stem, ""

        loader = factory(file_path)
        docs = loader.load()
        content = "\n\n".join(doc.page_content for doc in docs if doc.page_content)
        return Path(file_path).stem, content

    except ImportError:
        logger.warning(
            "langchain-community/unstructured not installed. "
            "Falling back to basic extraction for %s files. "
            "Install with: pip install langchain-community unstructured",
            ext,
        )
        return _extract_office_fallback(file_path, ext)
    except (OSError, ValueError, RuntimeError) as e:
        logger.warning("loader_dict extraction failed for %s: %s", file_path, e)
        return _extract_office_fallback(file_path, ext)


def _extract_office_fallback(file_path: str, ext: str) -> tuple[str, str]:
    """Basic fallback for Office formats when Unstructured is unavailable.

    Routes to dedicated extractors for each format.
    """
    if ext in ("doc", "docx"):
        return _extract_docx_fallback(file_path)
    if ext == "pptx":
        return _extract_pptx_fallback(file_path)
    if ext == "xlsx":
        return _extract_xlsx_fallback(file_path)
    if ext == "csv":
        return _extract_csv_fallback(file_path)
    return Path(file_path).stem, ""


def _extract_docx_fallback(file_path: str) -> tuple[str, str]:
    """DOCX fallback: extract paragraphs + tables via python-docx."""
    try:
        from docx import Document

        doc = Document(file_path)
        parts: list[str] = []

        # Extract paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Extract tables as markdown
        for table in doc.tables:
            md_table = _table_to_markdown(
                [[cell.text.strip() for cell in row.cells] for row in table.rows]
            )
            if md_table:
                parts.append(md_table)

        return Path(file_path).stem, "\n\n".join(parts)
    except ImportError:
        logger.debug("python-docx not installed for DOCX extraction of %s", file_path)
    except (OSError, ValueError, KeyError) as e:
        logger.warning("DOCX extraction failed for %s: %s", Path(file_path).name, e)

    return Path(file_path).stem, ""


def _extract_pptx_fallback(file_path: str) -> tuple[str, str]:
    """PPTX fallback: extract slide text + speaker notes via python-pptx."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = [f"## Slide {slide_num}"]

            # Shape text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_parts.append(text)

                # Tables inside slides
                if shape.has_table:
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    md_table = _table_to_markdown(rows)
                    if md_table:
                        slide_parts.append(md_table)

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"**Notes:** {notes}")

            parts.append("\n\n".join(slide_parts))

        return Path(file_path).stem, "\n\n---\n\n".join(parts)
    except ImportError:
        logger.debug("python-pptx not installed for PPTX extraction of %s", file_path)
    except (OSError, ValueError, KeyError) as e:
        logger.warning("PPTX extraction failed for %s: %s", Path(file_path).name, e)

    return Path(file_path).stem, ""


def _extract_xlsx_fallback(file_path: str) -> tuple[str, str]:
    """XLSX fallback: extract sheets as markdown tables via openpyxl."""
    try:
        import shutil
        import tempfile
        from openpyxl import load_workbook

        # openpyxl rejects non-.xlsx extensions — copy to a temp .xlsx if needed
        load_path = file_path
        tmp_path = None
        if not file_path.lower().endswith((".xlsx", ".xlsm", ".xltx", ".xltm")):
            tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
            tmp_path = tmp.name
            tmp.close()
            shutil.copy2(file_path, tmp_path)
            load_path = tmp_path

        try:
            wb = load_workbook(load_path, read_only=True, data_only=True)
            parts: list[str] = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[list[str]] = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(cell) if cell is not None else "" for cell in row])

                if not rows or all(all(c == "" for c in r) for r in rows):
                    continue

                md_table = _table_to_markdown(rows)
                if md_table:
                    header = f"## {sheet_name}" if len(wb.sheetnames) > 1 else ""
                    parts.append(f"{header}\n\n{md_table}" if header else md_table)

            wb.close()
            return Path(file_path).stem, "\n\n".join(parts)
        finally:
            if tmp_path:
                Path(tmp_path).unlink(missing_ok=True)
    except ImportError:
        logger.debug("openpyxl not installed for XLSX extraction of %s", file_path)
    except (OSError, ValueError, KeyError) as e:
        logger.warning("XLSX extraction failed for %s: %s", Path(file_path).name, e)
    except Exception as e:
        logger.warning("XLSX extraction unexpected error for %s: %s", Path(file_path).name, e)

    return Path(file_path).stem, ""


def _extract_csv_fallback(file_path: str) -> tuple[str, str]:
    """CSV fallback: parse with built-in csv module, render as markdown table."""
    try:
        with open(file_path, newline="", errors="ignore") as f:
            # Sniff dialect for robustness
            sample = f.read(8192)
            f.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                dialect = csv.excel

            reader = csv.reader(f, dialect)
            rows = [row for row in reader]

        if not rows:
            return Path(file_path).stem, ""

        md_table = _table_to_markdown(rows)
        return Path(file_path).stem, md_table
    except (OSError, csv.Error) as e:
        logger.warning("CSV extraction failed for %s: %s", Path(file_path).name, e)

    return Path(file_path).stem, ""


def _extract_pk_office_file(file_path: str) -> tuple[str, str]:
    """Attempt to extract content from a PK (ZIP) magic-byte file.

    Tries docx -> pptx -> xlsx in order, since we can't reliably determine
    which Office format it is without inspecting ZIP internals.
    """
    import zipfile

    # Peek inside the ZIP to guess the format from content_types or paths
    actual_format = None
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            names = set(zf.namelist())
            if any(n.startswith("word/") for n in names):
                actual_format = "docx"
            elif any(n.startswith("ppt/") for n in names):
                actual_format = "pptx"
            elif any(n.startswith("xl/") for n in names):
                actual_format = "xlsx"
    except (zipfile.BadZipFile, OSError):
        logger.debug("Could not inspect ZIP structure of %s", file_path)

    if actual_format == "docx":
        return _extract_docx_fallback(file_path)
    if actual_format == "pptx":
        return _extract_pptx_fallback(file_path)
    if actual_format == "xlsx":
        return _extract_xlsx_fallback(file_path)

    # Unknown ZIP — try each in order
    logger.debug("PK file %s: unknown internal structure, trying all Office parsers", Path(file_path).name)
    for extractor in (_extract_docx_fallback, _extract_pptx_fallback, _extract_xlsx_fallback):
        title, text = extractor(file_path)
        if text.strip():
            return title, text

    return Path(file_path).stem, ""


def _table_to_markdown(rows: list[list[str]]) -> str:
    """Convert a list of row-lists into a markdown table string.

    First row is treated as the header. Returns empty string for empty input.
    """
    if not rows:
        return ""

    # Ensure all rows have the same column count
    max_cols = max(len(r) for r in rows)
    normalized = [r + [""] * (max_cols - len(r)) for r in rows]

    header = normalized[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in normalized[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def _extract_plain_text(file_path: str) -> tuple[str, str]:
    """Direct text read for .txt, .md, .rst, and unknown files."""
    try:
        content = Path(file_path).read_text(errors="ignore")
        return Path(file_path).stem, content
    except OSError as e:
        logger.exception("Failed to read plain text file %s", file_path)
        return Path(file_path).stem, ""


def _get_marker_models() -> dict:
    """Lazy-load Marker models once for local PDF fallback path."""
    global _marker_models
    if _marker_models is None:
        from marker.models import create_model_dict

        _marker_models = create_model_dict()
    return _marker_models
