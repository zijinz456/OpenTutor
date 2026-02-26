"""Unified classification pipeline — 7-step ingestion.

Step 0: Preprocessing (SHA-256 dedup, filename regex, user preset)
Step 1: MIME detection (python-magic / filetype)
Step 2: Content extraction (Marker / python-pptx / trafilatura)
Step 3: LLM classification (content category)
Step 4: Course fuzzy matching (thefuzz)
Step 5: Store to ingestion_jobs table
Step 6: Dispatch to business tables (content_tree, assignments)

References:
- Papra: hash-during-stream + DB unique constraint for dedup
- Unstructured: partition() auto-routing
- Apache Tika: 5-layer MIME detection fallback
- thefuzz: fuzzy string matching for course names
"""

import hashlib
import logging
import mimetypes
import re
import uuid
import asyncio
from pathlib import Path

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from models.course import Course
from models.ingestion import IngestionJob, Assignment
from services.llm.router import get_llm_client

logger = logging.getLogger(__name__)

# ── Step 0: Filename regex patterns ──

FILENAME_PATTERNS = {
    r"(?i)lecture|slides|ppt|lec\d": "lecture_slides",
    r"(?i)chapter|textbook|reading|book": "textbook",
    r"(?i)hw|homework|assignment|problem.?set|ps\d": "assignment",
    r"(?i)exam|midterm|final|test|quiz": "exam_schedule",
    r"(?i)syllabus|schedule|outline": "syllabus",
    r"(?i)notes?|summary|review": "notes",
}


def classify_by_filename(filename: str) -> str | None:
    """Step 0: Try to classify by filename regex (zero LLM cost)."""
    for pattern, category in FILENAME_PATTERNS.items():
        if re.search(pattern, filename):
            return category
    return None


# ── Step 1: MIME detection ──

def detect_mime_type(filename: str, content_bytes: bytes | None = None) -> str:
    """Step 1: Detect MIME type. Falls back to extension-based detection."""
    # Try python-magic if available
    try:
        import magic
        if content_bytes:
            return magic.from_buffer(content_bytes[:8192], mime=True)
    except ImportError:
        pass

    # Fallback: extension-based
    mime, _ = mimetypes.guess_type(filename)
    return mime or "application/octet-stream"


# ── Step 2: Content extraction ──

async def extract_content(
    file_path: str | None,
    url: str | None,
    mime_type: str,
) -> str:
    """Step 2: Extract text content based on file type.

    Routes to appropriate extractor:
    - PDF → Marker
    - HTML/URL → trafilatura
    - PPTX → python-pptx
    - Images → multimodal LLM base64 (Phase 1+)
    - Plain text → direct read
    """
    if url and not file_path:
        # URL extraction (offload blocking I/O and enforce timeout)
        def _extract_url_text(target_url: str) -> str:
            import trafilatura
            from services.parser.url import extract_text_from_html

            downloaded = trafilatura.fetch_url(target_url)
            if not downloaded:
                return ""
            extracted = trafilatura.extract(downloaded, include_tables=True) or ""
            if extracted.strip():
                return extracted
            # Fallback for non-HTML or minimal pages (JSON/plain text).
            cleaned = extract_text_from_html(downloaded)
            return cleaned or (downloaded if isinstance(downloaded, str) else "")

        try:
            return await asyncio.wait_for(
                asyncio.to_thread(_extract_url_text, url),
                timeout=15,
            )
        except Exception as e:
            logger.warning(f"URL extraction failed: {e}")
            return ""

    if not file_path:
        return ""

    ext = Path(file_path).suffix.lower()

    if ext == ".pdf" or "pdf" in mime_type:
        return await _extract_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(file_path)
    elif ext in (".html", ".htm"):
        return _extract_html(file_path)
    elif ext in (".txt", ".md", ".rst"):
        return Path(file_path).read_text(errors="ignore")
    elif ext in (".docx",):
        return _extract_docx(file_path)
    else:
        # Try reading as text
        try:
            return Path(file_path).read_text(errors="ignore")
        except Exception:
            return ""


async def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF using Marker."""
    try:
        from marker.converters.pdf import PdfConverter
        from marker.models import create_model_dict

        models = create_model_dict()
        converter = PdfConverter(artifact_dict=models)
        rendered = converter(file_path)
        return rendered.markdown
    except ImportError:
        # Fallback: PyPDF2
        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            return ""
    except Exception as e:
        logger.warning(f"PDF extraction failed: {e}")
        return ""


def _extract_pptx(file_path: str) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts.append(f"## Slide {slide_num}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n\n".join(texts)
    except ImportError:
        logger.warning("python-pptx not installed, skipping PPTX extraction")
        return ""
    except Exception as e:
        logger.warning(f"PPTX extraction failed: {e}")
        return ""


def _extract_html(file_path: str) -> str:
    """Extract text from HTML using trafilatura."""
    try:
        import trafilatura
        with open(file_path) as f:
            html = f.read()
        return trafilatura.extract(html, include_tables=True) or ""
    except Exception as e:
        logger.warning(f"HTML extraction failed: {e}")
        return Path(file_path).read_text(errors="ignore")


def _extract_docx(file_path: str) -> str:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text)
    except ImportError:
        logger.warning("python-docx not installed, skipping DOCX extraction")
        return ""
    except Exception as e:
        logger.warning(f"DOCX extraction failed: {e}")
        return ""


# ── Step 3: LLM classification ──

CLASSIFICATION_PROMPT = """Classify this educational document into ONE of these categories:

Categories:
- lecture_slides: Lecture slides, presentations, class notes from professor
- textbook: Textbook chapters, readings, reference material
- assignment: Homework, problem sets, coding assignments
- exam_schedule: Exams, quizzes, midterms, finals, test dates
- syllabus: Course syllabus, schedule, grading policy
- notes: Student notes, study guides, review sheets
- other: Anything that doesn't fit above

Respond with ONLY the category name, nothing else.

Document (first 2000 chars):
{content}"""


async def classify_content(content: str, filename: str) -> str:
    """Step 3: LLM-based content classification.

    Only called when filename regex fails.
    Uses extract (lightweight) endpoint for cost efficiency.
    """
    # Try filename first (zero cost)
    category = classify_by_filename(filename)
    if category:
        return category

    if not content or len(content) < 50:
        return "other"

    try:
        client = get_llm_client()
        prompt = CLASSIFICATION_PROMPT.format(content=content[:2000])
        result = await client.extract(
            "You are a document classifier. Output only the category name.",
            prompt,
        )
        result = result.strip().lower()

        valid_categories = {
            "lecture_slides", "textbook", "assignment",
            "exam_schedule", "syllabus", "notes", "other",
        }
        if result in valid_categories:
            return result
        # Fuzzy match
        for cat in valid_categories:
            if cat in result:
                return cat
        return "other"
    except Exception as e:
        logger.warning(f"LLM classification failed: {e}")
        return "other"


# ── Step 4: Course fuzzy matching ──

async def match_course(
    db: AsyncSession,
    filename: str,
    content: str,
    user_id: uuid.UUID,
) -> uuid.UUID | None:
    """Step 4: Match ingested content to an existing course.

    Uses thefuzz for fuzzy string matching against course names.
    Falls back to None if no confident match (user assigns manually).
    """
    result = await db.execute(select(Course).where(Course.user_id == user_id))
    courses = result.scalars().all()

    if not courses:
        return None

    if len(courses) == 1:
        return courses[0].id

    # Try fuzzy matching
    try:
        from thefuzz import fuzz
    except ImportError:
        # Fallback: simple substring matching
        filename_lower = filename.lower()
        for course in courses:
            if course.name.lower() in filename_lower:
                return course.id
        return None

    best_score = 0
    best_match = None
    search_text = f"{filename} {content[:200]}".lower()

    for course in courses:
        score = fuzz.partial_ratio(course.name.lower(), search_text)
        if score > best_score:
            best_score = score
            best_match = course

    # Only match if confidence is high enough
    if best_score >= 70 and best_match:
        return best_match.id

    return None


# ── Step 5 & 6: Full pipeline ──

async def run_ingestion_pipeline(
    db: AsyncSession,
    user_id: uuid.UUID,
    file_path: str | None = None,
    url: str | None = None,
    filename: str = "",
    course_id: uuid.UUID | None = None,
    file_bytes: bytes | None = None,
) -> IngestionJob:
    """Run the full 7-step ingestion pipeline.

    Returns the IngestionJob with results.
    """
    # Step 0: SHA-256 dedup
    content_hash = None
    if file_bytes:
        content_hash = hashlib.sha256(file_bytes).hexdigest()
        # Check for duplicates
        existing = await db.execute(
            select(IngestionJob).where(
                IngestionJob.content_hash == content_hash,
                IngestionJob.user_id == user_id,
                IngestionJob.status == "completed",
            )
        )
        dupe = existing.scalar_one_or_none()
        if dupe:
            logger.info(f"Duplicate detected: {filename} (hash: {content_hash[:12]})")
            return dupe

    # Create job
    job = IngestionJob(
        user_id=user_id,
        source_type="file" if file_path else "url",
        original_filename=filename,
        url=url,
        content_hash=content_hash,
        course_id=course_id,
        course_preset=course_id is not None,
        status="extracting",
    )
    db.add(job)
    await db.flush()

    try:
        # Step 1: MIME detection
        if filename:
            job.mime_type = detect_mime_type(filename, file_bytes)

        # Step 2: Content extraction
        job.status = "extracting"
        extracted = await extract_content(file_path, url, job.mime_type or "")
        job.extracted_markdown = extracted

        if not extracted:
            job.status = "failed"
            job.error_message = "No content could be extracted"
            return job

        # Step 3: Classification
        job.status = "classifying"
        filename_category = classify_by_filename(filename)
        if filename_category:
            job.content_category = filename_category
            job.classification_method = "filename_regex"
        else:
            job.content_category = await classify_content(extracted, filename)
            job.classification_method = "llm_classification"

        # Step 4: Course matching (if not preset)
        if not course_id:
            matched_id = await match_course(db, filename, extracted, user_id)
            if matched_id:
                job.course_id = matched_id

        # Step 5: Status update
        job.status = "dispatching"

        # Step 6: Dispatch to business tables
        dispatch_result = await _dispatch_content(db, job)
        job.dispatched = True
        job.dispatched_to = dispatch_result
        job.status = "completed"

    except Exception as e:
        job.status = "failed"
        job.error_message = str(e)
        logger.error(f"Ingestion pipeline failed: {e}")

    await db.flush()
    return job


async def _dispatch_content(db: AsyncSession, job: IngestionJob) -> dict:
    """Step 6: Dispatch extracted content to appropriate business tables."""
    result = {}

    if not job.course_id or not job.extracted_markdown:
        return result

    category = job.content_category or "other"

    if category in ("lecture_slides", "textbook", "notes", "syllabus"):
        # Build content tree using PageIndex pattern
        from services.parser.pdf import _markdown_to_tree
        nodes = _markdown_to_tree(
            markdown=job.extracted_markdown,
            course_id=job.course_id,
            source_file=job.original_filename or "Untitled",
        )
        for node in nodes:
            # Normalize source metadata to the ingestion source (file/url).
            node.source_type = job.source_type
            node.source_file = job.original_filename
            db.add(node)
        result["content_tree"] = len(nodes)

    elif category == "assignment":
        # Extract assignment info
        assignment = Assignment(
            course_id=job.course_id,
            title=job.original_filename or "Assignment",
            description=job.extracted_markdown[:2000] if job.extracted_markdown else None,
            assignment_type="homework",
            source_ingestion_id=job.id,
        )
        db.add(assignment)
        result["assignments"] = 1

    elif category == "exam_schedule":
        # Extract exam dates using LLM
        assignment = Assignment(
            course_id=job.course_id,
            title=job.original_filename or "Exam",
            description=job.extracted_markdown[:2000] if job.extracted_markdown else None,
            assignment_type="exam",
            source_ingestion_id=job.id,
        )
        db.add(assignment)
        result["assignments"] = 1

    return result
